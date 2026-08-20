#!/usr/bin/env python
"""
cosmetics_check.py — presentation-quality sweep over the decks.

Only checks with a RELIABLE signal are included. An earlier version estimated text
overflow from shape geometry and produced ~1500 findings that were almost all noise
(it flagged one-word titles). That approach is gone.

  SHRUNK    PowerPoint itself recorded a normAutofit fontScale < 100%, i.e. it had to
            shrink the text to make it fit. That is overflow, on the authority of the
            renderer rather than a guess. <92% is visibly smaller than its neighbours.
  QUOTES    straight ' or " in prose, where these decks use curly typography
  ELLIPSIS  ".." (a typo for "…" or "...")
  EMPTY     a text shape holding only whitespace
"""
import glob
import os
import re

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
MONO = {"Courier New", "Consolas", "Menlo", "Monaco"}


def font_scale(shape):
    """Return PowerPoint's own autofit shrink factor, or None."""
    bodyPr = shape.text_frame._txBody.find(f"{A}bodyPr")
    if bodyPr is None:
        return None
    fit = bodyPr.find(f"{A}normAutofit")
    if fit is None:
        return None
    return int(fit.get("fontScale", 100000)) / 1000


def main():
    counts = {}
    for path in sorted(glob.glob(f"{ROOT}/*/slides/*.pptx")):
        base = os.path.basename(path)
        if base.startswith("~$") or "Understanding MCP" in base:
            continue
        out = []
        for i, slide in enumerate(Presentation(path).slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                full = tf.text
                if not full.strip():
                    if len(full):
                        out.append((f"EMPTY", i, f"shape {shape.shape_id} holds only whitespace"))
                    continue

                fs = font_scale(shape)
                if fs is not None and fs < 92:
                    out.append(("SHRUNK", i,
                                f"autofit shrank text to {fs:.0f}% :: {full[:58]!r}"))

                for p in tf.paragraphs:
                    mono = any((r.font.name or "") in MONO for r in p.runs)
                    t = "".join(r.text for r in p.runs)
                    if mono:
                        continue
                    if re.search(r"\.\.(?!\.)", t):
                        out.append(("ELLIPSIS", i, repr(t[:64])))
                    if re.search(r"[\"']", t):
                        out.append(("QUOTES", i, repr(t[:64])))

        if out:
            print(f"\n{'='*82}\n{base}\n{'='*82}")
            for kind, i, msg in out:
                print(f"  s{i:<3} {kind:<9} {msg}")
                counts[kind] = counts.get(kind, 0) + 1

    print("\n" + ", ".join(f"{v} {k}" for k, v in sorted(counts.items())) or "clean")


if __name__ == "__main__":
    main()
