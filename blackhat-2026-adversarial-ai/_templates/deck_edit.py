#!/usr/bin/env python
"""
deck_edit.py — targeted, verifiable edits applied directly to the .pptx decks.

The decks are the source of truth (the .md sources were retired 2026-08-06), so every
correction is applied here and recorded in EDITS below as an auditable list.

Text replacement preserves the run's formatting by editing the run in place. Where a
replacement spans runs, the paragraph is rebuilt using the first run's formatting.

    python _templates/deck_edit.py --dry-run
    python _templates/deck_edit.py
"""
import argparse
import copy
import glob
import os
import shutil

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# (deck glob, kind, payload)
#   replace : (old_substring, new_substring)  — applied to every matching paragraph
#   delete_slide : 1-based slide number       — applied after replacements
EDITS = [
    # ---- factual: "The Confused Deputy" (1988) has ONE author, Norm(an) Hardy.
    #      Verified 2026-08-06: ACM SIGOPS OSR 22(4), doi 10.1145/54289.871709.
    ("M5*", "replace",
     ("Hardy and Norman, the confused deputy (1988)",
      "Norm Hardy, “The Confused Deputy” (1988)")),

    # ---- instructor voice on a participant slide (M4 s10).
    ("M4*", "replace",
     ("SEC_ARTIFACT_VERIFICATION gates only load_artifact, which is used only in the "
      "instructor demo.",
      "SEC_ARTIFACT_VERIFICATION gates only load_artifact — the one code path that "
      "actually unpickles.")),
    ("M4*", "replace",
     ("Participants prove the finding statically. They never load the file.",
      "You prove the finding statically, from the bytecode. You never load the file.")),
    ("M4*", "replace", ("NOBODY RUNS THE RCE", "YOU PROVE IT WITHOUT RUNNING IT")),

    # ---- M7's "KNOWN GRADER BUG" slide needed deleting (instructor-voice, and the bug
    #      it described was fixed in eiger). KK removed it in PowerPoint on 2026-08-06,
    #      so no edit is required here. Left recorded so nobody re-adds it.

    # ---- Kicker labels addressed to the PRESENTER, on slides the room is reading.
    #      Rewritten audience-facing, keeping the same short all-caps rhythm.
    ("M2*", "replace", ("SAY IT OUT LOUD", "WHAT YOU JUST PROVED")),
    ("M2*", "replace", ("REFRAME IT", "WHY IT MATTERS")),
    ("M3*", "replace", ("START FROM WHY", "WHY TEAMS BUILD THIS")),
    ("M4*", "replace", ("PROJECT THE REAL OUTPUT", "THE REAL OUTPUT")),
    ("M4*", "replace", ("▶ DEMO / INSTRUCTOR ONLY", "▶ DEMO / WATCH ONLY")),
    ("M5*", "replace", ("DEMYSTIFY IT", "WHAT AN AGENT ACTUALLY IS")),
    ("M7*", "replace", ("GIVE THE DESIGN ITS DUE", "WHY THIS DESIGN EXISTS")),
    ("M7*", "replace", ("PRE-EMPT THE SHARP QUESTION", "THE SHARP QUESTION")),
    ("M8*", "replace", ("FRAME IT FAIRLY", "THE FAIR FRAMING")),
    ("M8*", "replace", ("SET EXPECTATIONS", "WHAT COUNTS AS SUCCESS")),

    # ---- Meta-commentary kickers: the slide referring to itself / to the deck.
    ("M4*", "replace", ("THE MOST USEFUL TABLE IN M4", "THE TABLE TO KEEP")),
    ("M5*", "replace", ("THE MOST IMPORTANT SLIDE IN M5", "THE POINT OF THIS MODULE")),
    ("M6*", "replace", ("THE KEY SLIDE", "THE IDEA THAT MATTERS")),

    # ---- "DISCUSSION BEAT" is a facilitation cue; M6 already uses plain "DISCUSSION".
    ("M2*", "replace", ("DISCUSSION BEAT", "DISCUSSION")),
    ("M3*", "replace", ("DISCUSSION BEAT", "DISCUSSION")),
    ("M5*", "replace", ("DISCUSSION BEAT", "DISCUSSION")),
    ("M7*", "replace", ("DISCUSSION BEAT", "DISCUSSION")),

    # ---- Opener s6: a single 341-character paragraph. Broken into three beats so it
    #      can be read from the back of the room.
    ("Opener*", "split_para",
     ("HubSpot, email, Slack, on Bedrock, with an autonomous planner→reviewer loop. At "
      "four points across the two days we do not just theorise — we drive Claude Code "
      "against Anna’s real code, live, and deliver a verdict: is this attack real, or "
      "theater? Everything you break in the lab, we test against a real app. That is "
      "the method you take home.",
      ["HubSpot, email, Slack, on Bedrock, with an autonomous planner→reviewer loop.",
       "At four points across the two days we don’t just theorise — we drive Claude "
       "Code against Anna’s real code, live, and deliver a verdict: is this attack "
       "real, or theater?",
       "Everything you break in the lab, we test against a real app. That is the "
       "method you take home."])),
]


def para_text(p):
    return "".join(r.text for r in p.runs)


def set_para(p, text):
    """Rewrite a paragraph to `text`, keeping the first run's formatting."""
    runs = p.runs
    if not runs:
        return
    rPr = runs[0]._r.find(f"{A}rPr")
    keep = copy.deepcopy(rPr) if rPr is not None else None
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    runs[0].text = text
    if keep is not None:
        cur = runs[0]._r.find(f"{A}rPr")
        if cur is not None:
            runs[0]._r.remove(cur)
        runs[0]._r.insert(0, keep)


def delete_slide(prs, index):
    """python-pptx has no public delete; drop the sldId and the relationship."""
    slides = prs.slides
    sldIdLst = slides._sldIdLst
    ids = list(sldIdLst)
    target = ids[index - 1]
    prs.part.drop_rel(target.rId)
    sldIdLst.remove(target)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    by_deck = {}
    for pattern, kind, payload in EDITS:
        for md in glob.glob(f"{ROOT}/*/slides/{pattern}.pptx"):
            if os.path.basename(md).startswith("~$") or "Understanding MCP" in md:
                continue
            by_deck.setdefault(md, []).append((kind, payload))

    total = 0
    for path, ops in sorted(by_deck.items()):
        prs = Presentation(path)
        name = os.path.basename(path)
        changed = 0
        print(f"\n{name}")

        for kind, payload in ops:
            if kind != "replace":
                continue
            old, new = payload
            hit = False
            for i, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for p in shape.text_frame.paragraphs:
                        t = para_text(p)
                        if old in t:
                            print(f"    s{i:<3} replace: {old[:58]!r}")
                            print(f"           ->      {new[:58]!r}")
                            if not args.dry_run:
                                set_para(p, t.replace(old, new))
                            changed += 1
                            hit = True
            if not hit:
                print(f"    !! NO MATCH for {old[:70]!r}")

        for kind, payload in ops:
            if kind != "split_para":
                continue
            old, parts = payload
            hit = False
            for i, slide in enumerate(prs.slides, 1):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for p in shape.text_frame.paragraphs:
                        if para_text(p).strip() != old.strip():
                            continue
                        print(f"    s{i:<3} split 1 paragraph ({len(old)} chars) "
                              f"into {len(parts)}")
                        hit = True
                        changed += 1
                        if args.dry_run:
                            continue
                        set_para(p, parts[0])
                        anchor = p._p
                        for extra in parts[1:]:
                            new_p = copy.deepcopy(p._p)
                            anchor.addnext(new_p)
                            anchor = new_p
                        # rewrite the clones' text
                        paras = shape.text_frame.paragraphs
                        idx = [q._p for q in paras].index(p._p)
                        for off, extra in enumerate(parts[1:], start=1):
                            set_para(paras[idx + off], extra)
            if not hit:
                print(f"    !! NO MATCH for split {old[:60]!r}")

        for kind, payload in ops:
            if kind != "delete_slide":
                continue
            n = payload
            title = ""
            for shape in prs.slides[n - 1].shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().split("\n")[0]
                    break
            print(f"    s{n:<3} DELETE SLIDE: {title[:60]!r}")
            if not args.dry_run:
                delete_slide(prs, n)
            changed += 1

        if changed and not args.dry_run:
            if not os.path.exists(path + ".prefix.bak"):
                shutil.copy2(path, path + ".prefix.bak")
            prs.save(path)
        total += changed

    print(f"\n{'DRY RUN — ' if args.dry_run else ''}{total} edits.")


if __name__ == "__main__":
    main()
