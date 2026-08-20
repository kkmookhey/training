#!/usr/bin/env python
"""
diff_deck_vs_source.py — what would a re-render change?

The .pptx decks were hand-edited in PowerPoint after their .md sources were written,
so the two have diverged. Re-rendering regenerates the deck FROM the .md, which means
any PowerPoint-only edit is silently destroyed.

Run this BEFORE any re-render.

    ONLY-IN-DECK  → PowerPoint-only content. A re-render DELETES this.
                    Backfill it into the .md first, or lose it.
    ONLY-IN-MD    → new source content not yet rendered. A re-render ADDS this.
                    (Expected for slides you just wrote.)
"""
import glob
import os
import re
import sys
from difflib import SequenceMatcher

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
MATCH = 0.62


def norm(s):
    return re.sub(r"[^a-z0-9]+", " ", s.lower()).strip()


def md_slides(path):
    src = re.sub(r"<!--.*?-->", "", open(path, encoding="utf-8").read(), flags=re.S)
    out = []
    for chunk in re.split(r"\n---\n", src):
        if not chunk.strip():
            continue
        head = next((l.lstrip("# ").strip()
                     for l in chunk.strip().split("\n") if l.startswith("#")), "")
        out.append((head, norm(chunk)))
    return out


def deck_slides(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
        title = (parts[1] if len(parts) > 1 else (parts[0] if parts else "")).strip()
        out.append((i, title.split("\n")[0][:58], norm(" ".join(parts))))
    return out


def main():
    total_lost = 0
    for md in sorted(glob.glob(f"{ROOT}/*/slides/*.md")):
        folder = os.path.dirname(md)
        decks = [p for p in glob.glob(f"{folder}/*.pptx")
                 if not os.path.basename(p).startswith("~$")
                 and "Understanding MCP" not in p]
        if not decks:
            continue
        pptx = max(decks, key=os.path.getsize) if len(decks) > 1 else decks[0]

        M = md_slides(md)
        D = deck_slides(pptx)

        # The renderer rewrites headings into punchy caps and condenses prose, so
        # character-level similarity badly under-matches. A rendered slide is
        # essentially a SUBSET of its source's words — so score by containment:
        # what fraction of the deck slide's words appear in the md slide.
        pairs = []
        for di, (n, title, dtext) in enumerate(D):
            dw = set(dtext.split())
            for mi, (head, mtext) in enumerate(M):
                if not dw or not mtext:
                    continue
                mw = set(mtext.split())
                pairs.append((len(dw & mw) / len(dw), di, mi))
        pairs.sort(reverse=True)
        ud, um = set(), set()
        for r, di, mi in pairs:
            if r < MATCH or di in ud or mi in um:
                continue
            ud.add(di)
            um.add(mi)

        lost = [D[i] for i in range(len(D)) if i not in ud and D[i][2]]
        added = [M[i] for i in range(len(M)) if i not in um and M[i][1]]

        if lost or added:
            print(f"\n{'='*76}\n{os.path.basename(pptx)}"
                  f"   ({len(D)} deck slides / {len(M)} md slides)\n{'='*76}")
        for n, title, _ in lost:
            print(f"  ONLY-IN-DECK  s{n:<3} {title!r}")
            total_lost += 1
        for head, _ in added:
            print(f"  only-in-md         {head[:58]!r}")

    print(f"\n{'='*76}")
    print(f"{total_lost} deck slides have no source counterpart — "
          f"a re-render would DELETE these.")
    print("Backfill them into the .md before rendering.")


if __name__ == "__main__":
    main()
