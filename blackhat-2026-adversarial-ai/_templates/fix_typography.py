#!/usr/bin/env python
"""
fix_typography.py — normalise straight quotes to curly in PROSE only.

M0, M2, M3, M4, M5, M7, M8 already use curly typography ("Eiger’s", "ISN’T").
M1 and M6 were rendered with straight quotes, so the set looks inconsistent on
screen. This converts prose quotes and leaves monospace runs alone — code must
keep its straight quotes or it stops being valid code.

    python _templates/fix_typography.py --dry-run
    python _templates/fix_typography.py
"""
import argparse
import glob
import os
import re
import shutil

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
MONO = {"Courier New", "Consolas", "Menlo", "Monaco"}

LSQUO, RSQUO, LDQUO, RDQUO = "‘", "’", "“", "”"


def curl(text):
    # Paired double quotes → “ … ” (non-greedy, must not span a run boundary).
    text = re.sub(r'"([^"]*)"', LDQUO + r"\1" + RDQUO, text)
    # Apostrophe inside a word: It's, don't, Iggy's, BCC'd, servers' → ’
    text = re.sub(r"(?<=\w)'(?=\w)", RSQUO, text)
    text = re.sub(r"(?<=s)'(?=\s|$)", RSQUO, text)
    return text


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    changed_total = 0
    for path in sorted(glob.glob(f"{ROOT}/*/slides/*.pptx")):
        base = os.path.basename(path)
        if base.startswith("~$") or "Understanding MCP" in base:
            continue
        prs = Presentation(path)
        hits = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    if any((r.font.name or "") in MONO for r in p.runs):
                        continue  # never touch code
                    for r in p.runs:
                        new = curl(r.text)
                        if new != r.text:
                            hits.append((i, r.text[:56], new[:56]))
                            if not args.dry_run:
                                r.text = new
        if hits:
            print(f"\n{base}")
            for i, old, new in hits:
                print(f"    s{i:<3} {old!r}")
                print(f"         -> {new!r}")
            changed_total += len(hits)
            if not args.dry_run:
                if not os.path.exists(path + ".typo.bak"):
                    shutil.copy2(path, path + ".typo.bak")
                prs.save(path)

    print(f"\n{'DRY RUN — ' if args.dry_run else ''}{changed_total} runs normalised.")


if __name__ == "__main__":
    main()
