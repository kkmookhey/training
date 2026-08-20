#!/usr/bin/env python
"""
fix_code_blocks.py — repair mangled monospace blocks in rendered slide decks.

WHY THIS EXISTS
---------------
The Claude Design render applies syntax highlighting by splitting each code line
into multiple runs (plain code + coloured comment/accent). It then places the
paragraph break BEFORE the coloured run instead of after it, so every highlighted
fragment is displaced onto the following line. On M5 s13 this pushed
`# VULNERABLE: allow everything` in front of a live `if` statement, i.e. the
Python on the slide was simply wrong.

This script re-imposes the .md source on the .pptx: one paragraph per source
line, byte-exact, with highlighting re-applied cleanly (comments keep the deck's
own comment colour). It is idempotent — re-run it after every render.

USAGE
    python _templates/fix_code_blocks.py --dry-run     # report only
    python _templates/fix_code_blocks.py               # write changes
    python _templates/fix_code_blocks.py --only M5     # one deck

Close the decks in PowerPoint first: an open file will overwrite these edits.
"""
import argparse
import copy
import glob
import os
import re
import shutil
import sys
from difflib import SequenceMatcher

from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

# Deck folder -> (source .md, rendered .pptx). Both discovered per module.
MONO_FONTS = {"Courier New", "Consolas", "Menlo", "Monaco", "Courier"}

# Languages where a trailing `#` / `//` is a comment we should re-colour.
COMMENT_MARKERS = {
    "python": "#", "py": "#", "bash": "#", "sh": "#", "shell": "#", "yaml": "#",
    "js": "//", "javascript": "//", "ts": "//", "typescript": "//", "java": "//",
}
# Fences we must never split (ASCII diagrams, JSON, plain output).
NO_SPLIT = {"", "text", "json", "http", "console", "diagram"}

MATCH_THRESHOLD = 0.55


def norm(s):
    """Aggressive normalisation for matching: the renderer transliterates
    box-drawing characters to ASCII, so compare on alphanumerics only."""
    return re.sub(r"[^a-z0-9]", "", s.lower())


def md_code_blocks(md_path):
    """Return [(lang, [lines])] for every fenced block outside speaker notes."""
    src = open(md_path, encoding="utf-8").read()
    src = re.sub(r"<!--.*?-->", "", src, flags=re.S)  # drop notes + header
    out = []
    for m in re.finditer(r"```([a-zA-Z]*)\n(.*?)```", src, flags=re.S):
        lang = m.group(1).lower()
        lines = m.group(2).rstrip("\n").split("\n")
        if len(lines) >= 2:
            out.append((lang, lines))
    return out


def code_shapes(prs):
    """Every shape that is predominantly monospace, in slide order."""
    found = []
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
            if len(runs) < 2:
                continue
            mono = sum(1 for r in runs if (r.font.name or "") in MONO_FONTS)
            if mono and mono >= len(runs) * 0.8:
                found.append((idx, shape))
    return found


def split_comment(line, marker):
    """Split a source line into (code, comment) at a real trailing comment.
    Quote-aware, and only counts a marker preceded by whitespace."""
    if not marker:
        return line, ""
    q = None
    i = 0
    while i < len(line):
        c = line[i]
        if q:
            if c == "\\":
                i += 2
                continue
            if c == q:
                q = None
        elif c in "\"'":
            q = c
        elif line.startswith(marker, i) and (i == 0 or line[i - 1] in " \t"):
            return line[:i], line[i:]
        i += 1
    return line, ""


def rpr_of(run):
    rPr = run._r.find(f"{A}rPr")
    return copy.deepcopy(rPr) if rPr is not None else None


def pick_formats(shape):
    """base rPr (most common) and comment rPr (a coloured run starting with a
    comment marker), sampled from what the renderer already produced."""
    runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
    counts = {}
    for r in runs:
        try:
            key = str(r.font.color.rgb)
        except Exception:
            key = "none"
        counts[key] = counts.get(key, 0) + 1
    base_key = max(counts, key=counts.get)

    base = comment = None
    for r in runs:
        try:
            key = str(r.font.color.rgb)
        except Exception:
            key = "none"
        if base is None and key == base_key:
            base = rpr_of(r)
        if comment is None and key != base_key and r.text.lstrip()[:2] in ("# ", "//"):
            comment = rpr_of(r)
    # Explicit None test: an lxml element with no children is falsy, so `or`
    # would silently discard a valid-but-empty comment rPr.
    return base, base if comment is None else comment


def rebuild(shape, lines, lang):
    """Replace the shape's text with `lines`, one paragraph per line."""
    base, comment = pick_formats(shape)
    marker = COMMENT_MARKERS.get(lang, "") if lang not in NO_SPLIT else ""
    tf = shape.text_frame
    txBody = tf._txBody

    for p in txBody.findall(f"{A}p"):
        txBody.remove(p)

    from pptx.oxml.ns import qn
    for line in lines:
        p = txBody.makeelement(qn("a:p"), {})
        txBody.append(p)
        code, cmt = split_comment(line, marker)
        for text, rPr in ((code, base), (cmt, comment)):
            if not text:
                continue
            r = p.makeelement(qn("a:r"), {})
            if rPr is not None:
                r.append(copy.deepcopy(rPr))
            t = p.makeelement(qn("a:t"), {})
            t.text = text
            r.append(t)
            p.append(r)


def process(md_path, pptx_path, dry_run):
    blocks = md_code_blocks(md_path)
    prs = Presentation(pptx_path)
    shapes = code_shapes(prs)
    if not blocks or not shapes:
        return 0, 0

    # Greedy best-match between source blocks and rendered shapes.
    pairs = []
    for bi, (lang, lines) in enumerate(blocks):
        bn = norm("\n".join(lines))
        for si, (slide_no, shape) in enumerate(shapes):
            sn = norm(shape.text_frame.text)
            if not sn or not bn:
                continue
            ratio = SequenceMatcher(None, bn, sn).ratio()
            pairs.append((ratio, bi, si))
    pairs.sort(reverse=True)

    used_b, used_s, chosen = set(), set(), []
    for ratio, bi, si in pairs:
        if ratio < MATCH_THRESHOLD or bi in used_b or si in used_s:
            continue
        used_b.add(bi)
        used_s.add(si)
        chosen.append((ratio, bi, si))

    fixed = 0
    for ratio, bi, si in sorted(chosen, key=lambda x: x[2]):
        lang, lines = blocks[bi]
        slide_no, shape = shapes[si]
        # "Already correct" must be judged LINE BY LINE. Comparing whole-shape
        # normalised text would pass a mangled block, because displacing a run
        # onto the next paragraph moves no characters — it only moves the breaks.
        # Compared EXACTLY, not normalised: the renderer also transliterates
        # box-drawing (└── ▼ ▶ →) to ASCII, which silently breaks the column
        # alignment of the ASCII diagrams. Source text is the authority.
        want = [l.rstrip() for l in lines]
        have = [("".join(r.text for r in p.runs)).rstrip()
                for p in shape.text_frame.paragraphs]
        if want == have:
            continue
        print(f"    s{slide_no:<3} match={ratio:.2f} lang={lang or '-':<6} "
              f"{len(lines)} lines  |  {lines[0][:52]!r}")
        if not dry_run:
            rebuild(shape, lines, lang)
        fixed += 1

    unmatched = [blocks[i][1][0][:50] for i in range(len(blocks)) if i not in used_b]
    for u in unmatched:
        print(f"    [no shape matched] {u!r}")

    if fixed and not dry_run:
        shutil.copy2(pptx_path, pptx_path + ".bak")
        prs.save(pptx_path)
    return fixed, len(unmatched)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dry-run", action="store_true")
    ap.add_argument("--only", default="")
    args = ap.parse_args()

    total_fixed = total_unmatched = 0
    for md in sorted(glob.glob(f"{ROOT}/*/slides/*.md")):
        folder = os.path.dirname(md)
        if args.only and args.only.lower() not in md.lower():
            continue
        decks = [p for p in glob.glob(f"{folder}/*.pptx")
                 if not os.path.basename(p).startswith("~$")
                 and "Understanding MCP" not in p]
        if not decks:
            continue
        pptx = max(decks, key=os.path.getsize) if len(decks) > 1 else decks[0]
        print(f"\n{os.path.relpath(pptx, ROOT)}")
        if os.path.exists(os.path.join(folder, "~$" + os.path.basename(pptx))):
            print("    !! OPEN IN POWERPOINT — close it and re-run. Skipped.")
            continue
        f, u = process(md, pptx, args.dry_run)
        total_fixed += f
        total_unmatched += u

    print(f"\n{'DRY RUN — ' if args.dry_run else ''}"
          f"{total_fixed} code blocks repaired, {total_unmatched} unmatched.")


if __name__ == "__main__":
    main()
